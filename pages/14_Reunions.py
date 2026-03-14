"""
Page 14 - Reunions de chantier
Preparation automatique de presentations PowerPoint pour les reunions.
Types: Reunion de chantier, reunion d'organisation, reunion interne, autre.
"""

import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import streamlit as st
import io
from datetime import datetime, date
from lib.helpers import page_setup, render_saas_sidebar, chantier_selector
from lib import db, storage
from utils import GLOBAL_CSS

page_setup(title="Reunions", icon="")
st.markdown(GLOBAL_CSS, unsafe_allow_html=True)

user_id = st.session_state.get("user_id")
render_saas_sidebar(user_id)

st.title("Reunions & Presentations")

chantier = chantier_selector(key="reunion_chantier")
if not chantier:
    st.stop()


# --- Types de reunions ---
REUNION_TYPES = {
    "Reunion de chantier": {
        "description": "Suivi d'avancement, points techniques, planning, securite",
        "sections": ["Ordre du jour", "Avancement des travaux", "Points techniques", "Planning", "Securite / Environnement", "Divers"],
    },
    "Reunion d'organisation": {
        "description": "Planification, coordination inter-lots, logistique",
        "sections": ["Ordre du jour", "Planning general", "Coordination inter-lots", "Logistique chantier", "Moyens humains et materiels", "Actions a mener"],
    },
    "Reunion interne": {
        "description": "Point d'equipe, objectifs, problemes internes",
        "sections": ["Ordre du jour", "Point equipe", "Objectifs", "Problemes rencontres", "Solutions proposees", "Actions"],
    },
    "Reunion client": {
        "description": "Presentation au maitre d'ouvrage, validation, decisions",
        "sections": ["Ordre du jour", "Etat d'avancement", "Budget", "Points de validation", "Decisions a prendre", "Prochaines etapes"],
    },
    "Autre": {
        "description": "Reunion personnalisee",
        "sections": ["Ordre du jour", "Points a discuter", "Decisions", "Actions"],
    },
}


def generate_pptx(reunion_data: dict) -> bytes:
    """Genere un PowerPoint professionnel pour la reunion."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Couleurs du theme
    PRIMARY = RGBColor(0x1E, 0x27, 0x61)      # Navy
    SECONDARY = RGBColor(0xCA, 0xDC, 0xFC)     # Ice blue
    ACCENT = RGBColor(0xF9, 0x61, 0x67)        # Coral
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
    TEXT_DARK = RGBColor(0x2D, 0x3A, 0x4A)
    TEXT_MUTED = RGBColor(0x6B, 0x7B, 0x8D)

    chantier_nom = reunion_data.get("chantier_nom", "Chantier")
    reunion_type = reunion_data.get("type", "Reunion de chantier")
    date_reunion = reunion_data.get("date", datetime.now().strftime("%d/%m/%Y"))
    lieu = reunion_data.get("lieu", "")
    participants = reunion_data.get("participants", [])
    sections_content = reunion_data.get("sections", {})
    notes_generales = reunion_data.get("notes_generales", "")
    numero_reunion = reunion_data.get("numero", "1")

    def add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape(slide, left, top, width, height, fill_color, alpha=None):
        from pptx.oxml.ns import qn
        shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        if alpha is not None:
            solidFill = shape.fill._fill
            srgb = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
            if srgb is not None:
                alphaElem = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha * 1000))})
                srgb.append(alphaElem)
        return shape

    def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    # ==================
    # SLIDE 1 : Titre
    # ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, PRIMARY)

    # Barre decorative laterale
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)

    # Titre principal
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 reunion_type.upper(), font_size=42, bold=True, color=WHITE, font_name="Calibri")

    # Nom du chantier
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
                 chantier_nom, font_size=28, bold=False, color=SECONDARY, font_name="Calibri")

    # Infos
    info_text = f"Reunion n {numero_reunion}  |  {date_reunion}"
    if lieu:
        info_text += f"  |  {lieu}"
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
                 info_text, font_size=16, bold=False, color=SECONDARY, font_name="Calibri")

    # ConducteurPro
    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                 "ConducteurPro - Gestion de chantier", font_size=12, bold=False, color=TEXT_MUTED, font_name="Calibri")

    # ==================
    # SLIDE 2 : Participants
    # ==================
    if participants:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, LIGHT_BG)

        # Barre top
        add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

        add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                     "PARTICIPANTS", font_size=32, bold=True, color=PRIMARY, font_name="Calibri")

        # Ligne separatrice
        add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT)

        # Tableau des participants
        from pptx.util import Emu
        cols = 3
        col_width = Inches(3.8)
        start_y = Inches(1.8)

        for idx, participant in enumerate(participants):
            col = idx % cols
            row = idx // cols
            x = Inches(0.8) + col * col_width
            y = start_y + Emu(int(row * Inches(0.6)))

            nom_p = participant.get("nom", participant) if isinstance(participant, dict) else participant
            role_p = participant.get("role", "") if isinstance(participant, dict) else ""

            display = f"{nom_p}"
            if role_p:
                display += f" - {role_p}"

            add_text_box(slide, x, y, col_width - Inches(0.3), Inches(0.5),
                         display, font_size=14, bold=False, color=TEXT_DARK, font_name="Calibri")

    # ==================
    # SLIDE 3 : Ordre du jour
    # ==================
    odj = sections_content.get("Ordre du jour", "")
    if odj:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, WHITE)

        add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

        add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                     "ORDRE DU JOUR", font_size=32, bold=True, color=PRIMARY, font_name="Calibri")

        add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT)

        # Points de l'ordre du jour
        points = [p.strip() for p in odj.split("\n") if p.strip()]
        y_pos = Inches(1.8)

        for i, point in enumerate(points):
            # Numero
            add_text_box(slide, Inches(0.8), y_pos, Inches(0.6), Inches(0.5),
                         f"{i+1}.", font_size=20, bold=True, color=ACCENT, font_name="Calibri")
            # Texte
            add_text_box(slide, Inches(1.5), y_pos, Inches(10), Inches(0.5),
                         point, font_size=16, bold=False, color=TEXT_DARK, font_name="Calibri")
            y_pos += Inches(0.55)

            if y_pos > Inches(6.5):
                break

    # ==================
    # SLIDES par section
    # ==================
    for section_name, section_text in sections_content.items():
        if section_name == "Ordre du jour" or not section_text.strip():
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, WHITE)

        add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

        add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                     section_name.upper(), font_size=32, bold=True, color=PRIMARY, font_name="Calibri")

        add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT)

        # Contenu de la section
        lines = [l.strip() for l in section_text.split("\n") if l.strip()]
        y_pos = Inches(1.8)

        for line in lines:
            # Detecter les titres (lignes commencant par # ou *)
            if line.startswith("#") or line.startswith("**"):
                clean = line.lstrip("#* ").rstrip("*")
                add_text_box(slide, Inches(0.8), y_pos, Inches(11), Inches(0.5),
                             clean, font_size=18, bold=True, color=PRIMARY, font_name="Calibri")
                y_pos += Inches(0.5)
            elif line.startswith("-") or line.startswith("*"):
                clean = line.lstrip("-* ")
                # Puce coloree
                add_text_box(slide, Inches(0.8), y_pos, Inches(0.4), Inches(0.4),
                             "●", font_size=10, bold=False, color=ACCENT, font_name="Calibri")
                add_text_box(slide, Inches(1.3), y_pos, Inches(10.5), Inches(0.4),
                             clean, font_size=14, bold=False, color=TEXT_DARK, font_name="Calibri")
                y_pos += Inches(0.45)
            else:
                add_text_box(slide, Inches(0.8), y_pos, Inches(11), Inches(0.4),
                             line, font_size=14, bold=False, color=TEXT_DARK, font_name="Calibri")
                y_pos += Inches(0.45)

            if y_pos > Inches(6.5):
                # Nouvelle slide si trop de contenu
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_bg(slide, WHITE)
                add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)
                add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                             f"{section_name.upper()} (suite)", font_size=32, bold=True, color=PRIMARY, font_name="Calibri")
                add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT)
                y_pos = Inches(1.8)

    # ==================
    # SLIDE : Notes generales
    # ==================
    if notes_generales and notes_generales.strip():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, LIGHT_BG)

        add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)

        add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                     "NOTES ET REMARQUES", font_size=32, bold=True, color=PRIMARY, font_name="Calibri")

        add_shape(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), ACCENT)

        add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                     notes_generales, font_size=14, bold=False, color=TEXT_DARK, font_name="Calibri")

    # ==================
    # SLIDE FINALE : Prochaine reunion
    # ==================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)

    add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)

    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "MERCI", font_size=48, bold=True, color=WHITE, font_name="Calibri",
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                 f"{reunion_type} - {chantier_nom}", font_size=20, bold=False, color=SECONDARY,
                 font_name="Calibri", alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                 f"Genere le {datetime.now().strftime('%d/%m/%Y a %H:%M')} par ConducteurPro",
                 font_size=12, bold=False, color=TEXT_MUTED, font_name="Calibri",
                 alignment=PP_ALIGN.CENTER)

    # Sauvegarder
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ==========================================
# Interface utilisateur
# ==========================================

tab_new, tab_history = st.tabs(["Nouvelle reunion", "Historique"])

with tab_new:
    st.subheader("Preparer une reunion")

    col1, col2 = st.columns(2)

    with col1:
        reunion_type = st.selectbox(
            "Type de reunion",
            list(REUNION_TYPES.keys()),
            key="reunion_type",
        )
        st.caption(REUNION_TYPES[reunion_type]["description"])

        date_reunion = st.date_input("Date de la reunion", value=date.today(), key="date_reunion")
        numero = st.number_input("Numero de la reunion", min_value=1, value=1, key="num_reunion")

    with col2:
        lieu = st.text_input("Lieu", placeholder="Ex: Salle de reunion du chantier", key="lieu_reunion")
        heure = st.time_input("Heure", key="heure_reunion")

    # Participants
    st.markdown("---")
    st.subheader("Participants")

    if "participants_list" not in st.session_state:
        st.session_state.participants_list = [{"nom": "", "role": ""}]

    for i, p in enumerate(st.session_state.participants_list):
        c1, c2, c3 = st.columns([3, 3, 1])
        st.session_state.participants_list[i]["nom"] = c1.text_input(
            "Nom", value=p["nom"], key=f"part_nom_{i}", label_visibility="collapsed",
            placeholder="Nom du participant"
        )
        st.session_state.participants_list[i]["role"] = c2.text_input(
            "Role", value=p["role"], key=f"part_role_{i}", label_visibility="collapsed",
            placeholder="Fonction / Entreprise"
        )
        if c3.button("X", key=f"del_part_{i}") and len(st.session_state.participants_list) > 1:
            st.session_state.participants_list.pop(i)
            st.rerun()

    if st.button("+ Ajouter un participant"):
        st.session_state.participants_list.append({"nom": "", "role": ""})
        st.rerun()

    # Sections de la reunion
    st.markdown("---")
    st.subheader("Contenu de la presentation")

    sections = REUNION_TYPES[reunion_type]["sections"]
    sections_content = {}

    for section in sections:
        sections_content[section] = st.text_area(
            section,
            key=f"section_{section}",
            height=120,
            placeholder=f"Saisissez les points pour la section '{section}'...\nUtilisez des tirets (-) pour les listes\nUtilisez ** pour les sous-titres",
            help="Chaque ligne = un point. Commencez par - pour une puce, ** pour un sous-titre."
        )

    # Notes generales
    notes = st.text_area(
        "Notes et remarques generales",
        key="notes_reunion",
        height=80,
        placeholder="Informations complementaires, rappels, etc.",
    )

    # Documents du chantier a joindre
    st.markdown("---")
    st.subheader("Documents de reference")
    docs = db.get_documents(user_id=user_id, chantier_id=chantier["id"])
    if docs:
        doc_names = [d.get("nom", "N/A") for d in docs]
        selected_docs = st.multiselect(
            "Selectionner les documents a mentionner dans la reunion",
            doc_names,
            key="docs_reunion",
        )
    else:
        selected_docs = []
        st.info("Aucun document disponible pour ce chantier. Uploadez des documents dans la section Documents.")

    # Generation
    st.markdown("---")

    if st.button("Generer la presentation PowerPoint", type="primary", width="stretch"):
        if not any(v.strip() for v in sections_content.values()):
            st.error("Veuillez remplir au moins une section avant de generer la presentation.")
        else:
            with st.spinner("Generation de la presentation en cours..."):
                # Preparer les donnees
                participants_clean = [
                    p for p in st.session_state.participants_list
                    if p.get("nom", "").strip()
                ]

                reunion_data = {
                    "chantier_nom": chantier.get("nom", chantier.get("client_nom", "Chantier")),
                    "type": reunion_type,
                    "date": date_reunion.strftime("%d/%m/%Y"),
                    "heure": heure.strftime("%H:%M") if heure else "",
                    "lieu": lieu,
                    "numero": str(numero),
                    "participants": participants_clean,
                    "sections": sections_content,
                    "notes_generales": notes,
                    "documents_ref": selected_docs,
                }

                try:
                    pptx_bytes = generate_pptx(reunion_data)

                    # Nom du fichier
                    date_str = date_reunion.strftime("%Y%m%d")
                    safe_type = reunion_type.replace(" ", "_").replace("'", "")
                    filename = f"Reunion_{safe_type}_{date_str}_n{numero}.pptx"

                    # Telecharger directement
                    st.download_button(
                        label="Telecharger la presentation",
                        data=pptx_bytes,
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary",
                    )

                    # Sauvegarder dans Supabase
                    try:
                        doc_result = storage.upload_generated_document(
                            file_bytes=pptx_bytes,
                            filename=filename,
                            chantier_id=chantier["id"],
                            famille="PV",
                            doc_type="Reunion",
                            metadata={
                                "reunion_type": reunion_type,
                                "date_reunion": date_reunion.isoformat(),
                                "numero": str(numero),
                                "participants": [p.get("nom", "") for p in participants_clean],
                            },
                        )
                        if doc_result:
                            st.success(f"Presentation '{filename}' generee et sauvegardee !")
                        else:
                            st.warning("Presentation generee mais non sauvegardee dans le stockage.")
                    except Exception:
                        st.info("Presentation generee. Sauvegarde dans le stockage non disponible.")

                    # Sauvegarder la reunion en DB
                    try:
                        db.save_reunion(user_id, chantier["id"], {
                            "type": reunion_type,
                            "date_reunion": date_reunion.isoformat(),
                            "numero": int(numero),
                            "lieu": lieu,
                            "participants": participants_clean,
                            "sections": sections_content,
                            "notes": notes,
                            "filename": filename,
                        })
                    except Exception:
                        pass

                except ImportError:
                    st.error("La bibliotheque python-pptx n'est pas installee. Contactez l'administrateur.")
                    st.code("pip install python-pptx", language="bash")
                except Exception as e:
                    st.error(f"Erreur lors de la generation : {e}")


with tab_history:
    st.subheader("Reunions precedentes")

    try:
        reunions = db.get_reunions(user_id=user_id, chantier_id=chantier["id"])
        if reunions:
            for r in reunions:
                r_type = r.get("type", "Reunion")
                r_date = str(r.get("date_reunion", ""))[:10]
                r_num = r.get("numero", "")
                r_filename = r.get("filename", "")
                r_participants = r.get("participants", [])

                with st.expander(f"Reunion n{r_num} - {r_type} ({r_date})", expanded=False):
                    if r.get("lieu"):
                        st.caption(f"Lieu : {r['lieu']}")
                    if r_participants:
                        noms = [p.get("nom", p) if isinstance(p, dict) else str(p) for p in r_participants]
                        st.caption(f"Participants : {', '.join(noms)}")
                    if r.get("sections"):
                        for sec_name, sec_text in r["sections"].items():
                            if sec_text and sec_text.strip():
                                st.markdown(f"**{sec_name}**")
                                st.text(sec_text)
                    if r.get("notes") and r["notes"].strip():
                        st.markdown("**Notes**")
                        st.text(r["notes"])
        else:
            st.info("Aucune reunion enregistree pour ce chantier.")
    except Exception:
        st.info("Historique des reunions non disponible. La table 'reunions' n'existe peut-etre pas encore dans la base de donnees.")
